# streamlit_app.py
# One-input Streamlit app → generates a left-rail PPTX slide from the customer name.
# No runtime pip/subprocess. Supports OpenAI/Anthropic (optional) or No-LLM presets.

import io
import json
from datetime import datetime

import streamlit as st

# ---- Required deps (no runtime install) ----
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    st.error(
        "Missing dependency: **python-pptx**.\n\n"
        "Add `python-pptx` to requirements.txt, redeploy, and rerun the app."
    )
    st.stop()

# Optional (for calling OpenAI/Anthropic). App still works without it in No-LLM mode.
try:
    import requests
except ImportError:
    requests = None

# ---------------------------- UI ----------------------------
st.set_page_config(page_title="Customer Update Slide Generator", page_icon="📊", layout="centered")
st.title("📊 One-Input Customer Update Slide (PPTX)")

with st.expander("What this does", expanded=False):
    st.markdown(
        """
        - Enter a **customer name** (e.g., *Rugs USA*).  
        - Choose **OpenAI**, **Anthropic**, or **No-LLM** (presets/smart defaults).  
        - Click **Generate PPTX** to download your left-rail slide.  
        - No runtime installs; dependencies come from requirements.txt.
        """
    )

customer = st.text_input("Customer name", "Rugs USA")
c1, c2 = st.columns(2)
with c1:
    accent_hex = st.color_picker("Accent color", value="#08574A")
with c2:
    logo_file = st.file_uploader("Optional: Upload customer logo (PNG)", type=["png"])

provider_options = ["No-LLM (preset/smart template)"]
if requests is not None:
    provider_options = ["OpenAI", "Anthropic"] + provider_options
provider = st.selectbox("Generator", provider_options)

api_key = ""
if provider in ("OpenAI", "Anthropic"):
    api_key = st.text_input(f"{provider} API key", type="password", placeholder="Paste key here")

st.divider()

# ---------------------------- Content generation ----------------------------
SYSTEM_INSTRUCTIONS = """You are a precise enterprise research/writing assistant for Customer Success Managers.
Generate concise, presentation-ready content based on broadly available public context for the named company.
Keep bullets short (6–14 words). Avoid speculation or marketing fluff. If specifics are unclear, give conservative,
generically relevant bullets appropriate to the company's business model."""

USER_PROMPT_TMPL = """Company: {customer}

Task: Draft sections for an executive slide titled "Customer Updates – Strategic Priorities & Supply Chain Context".
Keep it concise, factual, and presentation-ready. Use bullets where applicable. No fluff.

Return strict JSON with keys:
- corporate_vision: string (2–3 sentences).
- business_strategies: array of 4–6 short bullets.
- supply_chain_contribution: array of 4–6 short bullets.
- risks_of_supply_chain_failure: array of 4–6 short bullets.
- critical_capabilities: array of 4–6 short bullets.

Rules:
- Base on widely reported info (brand positioning, retail/e-commerce ops, logistics practices).
- If specifics are not well-documented, provide conservative, relevant bullets for the company type.
- Do not include sources or URLs. Return only JSON.
"""

def rgb_hex_to_tuple(hex_color: str):
    h = hex_color.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def call_openai_chat(key: str, prompt: str) -> str:
    if requests is None:
        raise RuntimeError("The 'requests' package is required for OpenAI mode.")
    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    data = {
        "model": "gpt-4o-mini",
        "temperature": 0.2,
        "max_tokens": 900,
        "messages": [
            {"role": "system", "content": SYSTEM_INSTRUCTIONS},
            {"role": "user", "content": prompt},
        ],
    }
    r = requests.post(url, headers=headers, json=data, timeout=60)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"]

def call_anthropic(key: str, prompt: str) -> str:
    if requests is None:
        raise RuntimeError("The 'requests' package is required for Anthropic mode.")
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        "x-api-key": key,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json",
    }
    data = {
        "model": "claude-3-5-sonnet-20240620",
        "max_tokens": 900,
        "temperature": 0.2,
        "system": SYSTEM_INSTRUCTIONS,
        "messages": [{"role": "user", "content": prompt}],
    }
    r = requests.post(url, headers=headers, json=data, timeout=60)
    r.raise_for_status()
    blocks = r.json().get("content", [])
    return "".join(b.get("text", "") for b in blocks if b.get("type") == "text")

def safe_parse_json(text: str):
    try:
        return json.loads(text)
    except Exception:
        pass
    import re
    blocks = re.findall(r"```(?:json)?\s*(\{.*?\})\s*```", text, flags=re.S)
    for b in blocks:
        try:
            return json.loads(b)
        except Exception:
            continue
    raise ValueError("Could not parse JSON from model output.")

# Preset(s) for No-LLM mode (you can add more customers here)
PRESETS = {
    "Rugs USA": {
        "corporate_vision":
            "Rugs USA helps customers turn houses into homes by offering a curated, "
            "stylish, high-quality assortment at compelling prices, delivered through a "
            "seamless digital shopping and support experience.",
        "business_strategies": [
            "Expand ‘house of brands’ to broaden assortment and reach.",
            "Accelerate operational efficiency via DC consolidation and lean processes.",
            "Win digital: strengthen DTC and marketplaces for scaled demand.",
            "Refresh catalog with designer and partner collaborations.",
            "Tighten pricing, merchandising, and returns to lift unit economics.",
        ],
        "supply_chain_contribution": [
            "Faster delivery through optimized DC network and slotting.",
            "Flexible routing to serve DTC site and marketplaces.",
            "Rapid vendor onboarding with packaging and QC compliance.",
            "Inventory visibility and demand-aligned PO planning.",
            "Carrier and 3PL partnerships for cost and reliability.",
        ],
        "risks_of_supply_chain_failure": [
            "Debt pressure constraining capex for ops improvements.",
            "Housing-linked demand volatility driving forecast error.",
            "WMS/OMS outages disrupting pick-pack-ship rhythm.",
            "Integration slippage from M&A or DC consolidation.",
            "Supplier or carrier delays impacting service levels.",
        ],
        "critical_capabilities": [
            "Modern fulfillment stack: SaaS WMS/OMS and mobile tools.",
            "SKU-level forecasting and returns analytics.",
            "Marketplace and carrier relationship management.",
            "Integration playbooks for synergy capture and cost control.",
            "Governance: change control, DR, and vendor SLAs.",
        ],
    }
}

def generate_sections(customer_name: str, mode: str, key: str):
    # No-LLM preset / generic template
    if mode == "No-LLM (preset/smart template)":
        if customer_name in PRESETS:
            return PRESETS[customer_name]
        # generic fallback for unknown brands
        return {
            "corporate_vision":
                f"{customer_name} delivers curated, great-value products through a digital-first, "
                "customer-centric experience focused on quality, style, and convenience.",
            "business_strategies": [
                "Broaden assortment and strengthen brand positioning.",
                "Improve fulfillment speed and delivery reliability.",
                "Scale profitability via pricing and returns discipline.",
                "Grow demand through marketplaces and partnerships.",
                "Invest in data-driven merchandising and CX.",
            ],
            "supply_chain_contribution": [
                "Optimize DC footprint and pick-pack efficiency.",
                "Enable omnichannel routing and inventory visibility.",
                "Tighten supplier onboarding and packaging compliance.",
                "Leverage carrier mix for cost and service balance.",
                "Deploy WMS/OMS dashboards for real-time control.",
            ],
            "risks_of_supply_chain_failure": [
                "Forecast error from macro or seasonal volatility.",
                "System outages disrupting order flow or labeling.",
                "Carrier delays or capacity constraints.",
                "Supplier quality or lead-time variability.",
                "Cost inflation eroding contribution margins.",
            ],
            "critical_capabilities": [
                "Reliable WMS/OMS with strict change control.",
                "SKU-level demand and returns analytics.",
                "S&OP cadence linking demand to supply plans.",
                "Partner SLAs and performance governance.",
                "Contingency playbooks and DR testing.",
            ],
        }

    # OpenAI
    if mode == "OpenAI":
        if not requests:
            raise RuntimeError("`requests` package not available.")
        if not key:
            raise RuntimeError("Missing OpenAI API key.")
        text = call_openai_chat(key, USER_PROMPT_TMPL.format(customer=customer_name))
        return safe_parse_json(text)

    # Anthropic
    if mode == "Anthropic":
        if not requests:
            raise RuntimeError("`requests` package not available.")
        if not key:
            raise RuntimeError("Missing Anthropic API key.")
        text = call_anthropic(key, USER_PROMPT_TMPL.format(customer=customer_name))
        return safe_parse_json(text)

    raise RuntimeError("Unknown provider.")

# ---------------------------- PPTX rendering ----------------------------
def add_logo(slide, file_uploader_obj, width_in=1.6):
    try:
        if file_uploader_obj is None:
            return
        data = file_uploader_obj.getvalue()
        slide.shapes.add_picture(io.BytesIO(data), Inches(11.4), Inches(0.3), width=Inches(width_in))
    except Exception:
        pass

def add_rule(slide, y_in, color, thickness=2.2):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_in), Inches(12.3), Inches(0.02*thickness))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def left_label(slide, text, y_in):
    box = slide.shapes.add_textbox(Inches(0.3), Inches(y_in), Inches(2.5), Inches(1.0))
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = RGBColor(0,0,0)

def right_title_body(slide, title, body, y_in):
    tbox = slide.shapes.add_textbox(Inches(3.0), Inches(y_in), Inches(9.6), Inches(0.7))
    tf = tbox.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(22); p.font.bold = True

    bbox = slide.shapes.add_textbox(Inches(3.0), Inches(y_in+0.55), Inches(9.6), Inches(1.4))
    btf = bbox.text_frame; btf.clear()
    if isinstance(body, str):
        btf.paragraphs[0].text = body
        btf.paragraphs[0].font.size = Pt(14)
    else:
        for i, line in enumerate(body):
            pp = btf.add_paragraph() if i else btf.paragraphs[0]
            pp.text = line
            pp.font.size = Pt(14)
            pp.space_after = Pt(2)

def build_slide(customer_name: str, accent_rgb, sections: dict, logo_bytes=None):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Customer Updates"
    p.font.size = Pt(44); p.font.bold = True

    # Logo
    add_logo(slide, logo_file)

    # Accent vertical bar
    vbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(1.2), Inches(0.06), Inches(5.8))
    vbar.fill.solid(); vbar.fill.fore_color.rgb = RGBColor(*accent_rgb)
    vbar.line.fill.background()

    # Rows
    y, row_h = 1.25, 1.4
    rule_color = RGBColor(165,175,190)

    rows = [
        ("Corporate\nVision", "Corporate Vision", sections["corporate_vision"]),
        ("Business\nStrategies", "Business Strategies", ["• " + s for s in sections["business_strategies"]]),
        ("Supply Chain\nContribution", "Supply Chain Contribution", ["• " + s for s in sections["supply_chain_contribution"]]),
        ("Risks of\nSupply Chain\nFailure", "Risks of Supply Chain Failure", ["• " + s for s in sections["risks_of_supply_chain_failure"]]),
        ("Critical\nCapabilities", "Critical Capabilities", ["• " + s for s in sections["critical_capabilities"]]),
    ]
    for rail, title, body in rows:
        left_label(slide, rail, y - 0.05)
        right_title_body(slide, title, body, y - 0.02)
        add_rule(slide, y + row_h, rule_color)
        y += row_h

    # Footer
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4))
    ft = foot.text_frame; ft.clear()
    fp = ft.paragraphs[0]
    fp.text = f"{customer_name}  •  {datetime.now().strftime('%b %d, %Y')}"
    fp.font.size = Pt(12); fp.font.color.rgb = RGBColor(90,90,90)
    fp.alignment = PP_ALIGN.RIGHT

    return prs

# ---------------------------- Generate ----------------------------
if st.button("Generate PPTX"):
    try:
        # Build content
        sections = generate_sections(customer.strip(), provider, api_key.strip())
        # Render pptx
        prs = build_slide(customer.strip(), rgb_hex_to_tuple(accent_hex), sections, logo_file)
        buf = io.BytesIO()
        prs.save(buf); buf.seek(0)
        st.success("PPTX generated successfully!")
        st.download_button(
            "⬇️ Download PPTX",
            data=buf,
            file_name=f"{customer.strip().replace(' ','_')}_Customer_Updates_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"Generation failed: {e}")

st.caption("Use **No-LLM** for presets, or add an API key in the UI for tailored content.")
